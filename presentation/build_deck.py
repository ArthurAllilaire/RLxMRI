#!/usr/bin/env python3
"""Generate the RLxMRI FYP presentation deck from the agreed script.

Produces an editable .pptx: titles, on-slide bullets, embedded report
figures (where they exist) or labelled placeholders for figures still to be
made, and the full speaker script in each slide's notes pane.

Run:  python presentation/build_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = os.path.join(REPO, "report_latex", "imgs")
FIGS = os.path.join(REPO, "presentation", "figs")   # purpose-built talk figures
OUT = os.path.join(REPO, "presentation", "RLxMRI_presentation.pptx")


def resolve(p):
    """Find an image by name in the talk-figs dir first, then report imgs."""
    for base in (FIGS, IMG):
        full = os.path.join(base, p)
        if os.path.exists(full):
            return full
    return None

# Imperial brand colours (report uses ImperialBlue = RGB 0,0,205)
IMPERIAL = RGBColor(0, 0, 205)
NAVY = RGBColor(0, 30, 90)
INK = RGBColor(30, 30, 35)
GREY = RGBColor(120, 120, 128)
PLACEHOLDER_BG = RGBColor(238, 240, 248)
PLACEHOLDER_LINE = RGBColor(170, 178, 205)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def title_band(slide, title, kicker=None):
    """Coloured title strip at the top."""
    band = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SW, Inches(1.15))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    band.shadow.inherit = False
    tf = band.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def bullets_box(slide, items, left, top, width, height, size=16):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (lvl, text, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = ("• " if lvl == 0 else "– ") + text
        run.font.size = Pt(size - lvl * 2)
        run.font.color.rgb = IMPERIAL if bold else INK
        run.font.bold = bold
    return tb


def _fit_centered(slide, full, left, top, box_w, box_h):
    """Place an image scaled to fit (box_w, box_h), preserving aspect, centred."""
    iw, ih = Image.open(full).size
    scale = min(box_w / iw, box_h / ih)
    w = Emu(int(iw * scale))
    h = Emu(int(ih * scale))
    x = Emu(int(left + (box_w - w) / 2))
    y = Emu(int(top + (box_h - h) / 2))
    slide.shapes.add_picture(full, x, y, width=w, height=h)


def image_or_placeholder(slide, paths, left, top, width, height, label):
    """Embed first existing image(s) fitted into the box, else a placeholder."""
    existing = [resolve(p) for p in paths if resolve(p)]
    if existing:
        # stack up to two images vertically, each fitted into its sub-box
        n = len(existing[:2])
        gap = Inches(0.12)
        each_h = Emu(int((height - gap * (n - 1)) / n))
        y = top
        for full in existing[:2]:
            _fit_centered(slide, full, left, y, width, each_h)
            y = Emu(y + each_h + gap)
        return
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = PLACEHOLDER_BG
    box.line.color.rgb = PLACEHOLDER_LINE
    box.line.dash_style = 2
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY
    p.font.bold = True


# ---------------------------------------------------------------------------
# SLIDE 1 — Title
# ---------------------------------------------------------------------------
s = add_slide()
band = s.shapes.add_shape(1, Inches(0), Inches(2.4), SW, Inches(2.7))
band.fill.solid(); band.fill.fore_color.rgb = NAVY
band.line.fill.background(); band.shadow.inherit = False
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.1))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Adaptive Quantitative MRI Sequence Design"
p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)
p2 = tf.add_paragraph(); p2.text = "with Reinforcement Learning"
p2.font.size = Pt(40); p2.font.bold = True; p2.font.color.rgb = RGBColor(255, 255, 255)
sub = s.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.5))
stf = sub.text_frame; stf.word_wrap = True
for i, line in enumerate([
        "Arthur Allilaire  ·  Imperial College London  ·  BEng Final-Year Project",
        "Supervisors: Andreas Wetscherek (ICR)  ·  Wayne Luk (Imperial)"]):
    pp = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
    pp.text = line; pp.font.size = Pt(16); pp.font.color.rgb = INK
hook = s.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6))
hp = hook.text_frame.paragraphs[0]
hp.text = "Can an agent learn to design MRI scans on the fly — better than a fixed protocol?"
hp.font.size = Pt(15); hp.font.italic = True; hp.font.color.rgb = IMPERIAL
set_notes(s, "Title. One line: can an agent learn to design MRI scans on the fly, "
             "better than a fixed protocol? Name, supervisors, BEng FYP.")

# ---------------------------------------------------------------------------
# Helper for standard content slides
# ---------------------------------------------------------------------------
def bullets_as_text(items):
    """Flatten a (level, text, bold) bullet list into indented plain text."""
    lines = []
    for lvl, text, _bold in items:
        prefix = ("    " * lvl) + ("• " if lvl == 0 else "– ")
        lines.append(prefix + text)
    return "\n".join(lines)


def content_slide(title, kicker, bullets, images, img_label, notes="",
                  text_w=4.5, has_img=True, full_bullets=None):
    s = add_slide()
    title_band(s, title)            # kicker no longer printed on the slide
    if has_img:
        bullets_box(s, bullets, Inches(0.5), Inches(1.45),
                    Inches(text_w), Inches(5.6))
        image_or_placeholder(s, images, Inches(text_w + 0.7), Inches(1.5),
                             Inches(13.333 - text_w - 1.2), Inches(5.4),
                             img_label)
    else:
        bullets_box(s, bullets, Inches(0.5), Inches(1.45),
                    Inches(12.3), Inches(5.6))
    # Notes hold the marking-band tag, the spoken script, and the fuller bullet
    # version, so the on-slide bullets can stay terse without losing detail.
    parts = []
    if kicker:
        parts.append("[" + kicker + "]")
    if notes:
        parts.append(notes)
    if full_bullets:
        parts.append("— Slide bullets (detail) —\n" + bullets_as_text(full_bullets))
    set_notes(s, "\n\n".join(parts))
    return s


# SLIDE 2 — Why it matters
content_slide(
    "Why this matters: the scan-time bottleneck",
    "Framing (15%) — motivate the problem from a real clinical pressure",
    [(0, "MR-Linac: scan time is the bottleneck", True),
     (0, "qMRI: measure T1/T2 — many timed images + a fit", False),
     (0, "Today: one FIXED protocol for every patient", True),
     (0, "Best timing depends on the unknown tissue", True)],
    ["FIG1_clinical.png"],
    "FIG-1 (make / source)\nClean MR-Linac or qMRI T1-map image",
    "Clinical hook. MRI gives soft-tissue contrast that guides radiotherapy; acquisition "
    "time is the bottleneck. qMRI measures T1/T2 by acquiring several timed images and "
    "fitting a relaxation curve — more scan time. Today timings are fixed for every "
    "patient, but the most informative timing depends on the unknown tissue. That tension "
    "is the problem.",
    full_bullets=[
        (0, "MRI guides radiotherapy (MR-Linac) — but acquisition time is the clinical bottleneck", True),
        (1, "longer scans → motion artefacts, less on-table adaptation", False),
        (0, "Quantitative MRI (qMRI) measures physical tissue constants T1, T2", True),
        (1, "needs several images at different timings, then a curve fit → even more time", False),
        (0, "Today those timings are a FIXED protocol — same for every patient", True),
        (0, "But the most informative timing depends on the tissue you're measuring…", True),
        (1, "…which is exactly what you don't know yet", False)])

# SLIDE 3 — Gap + thesis
content_slide(
    "Adaptive qMRI: the idea, the gap, and our thesis",
    "Framing (15%) — state the gap as a TENSION, then the novelty claim",
    [(0, "Adaptive MRI ≈ 2.5× faster (Beracha)", False),
     (0, "Adaptive qMRI: Bayesian model only", False),
     (0, "RL in MRI: non-quantitative goals only", False),
     (0, "Thesis: first RL agent for adaptive qMRI", True),
     (0, "phantom (A1) → simulator (A2) → RL (A3)", False)],
    ["fig2_capability.png"],
    "FIG-2 capability matrix\n(works × learned/adaptive/\nquantitative)",
    "Adaptive idea: choose next acquisition from current estimate (Beracha ~2.5x). The gap "
    "as a tension: adaptive qMRI is Bayesian-model-only; RL in MRI targets non-quantitative "
    "goals; nobody combined learned + adaptive + quantitative. Thesis: first RL agent for "
    "adaptive qMRI, conditioned on fitted T1, scored on T1 error. Roadmap to the three "
    "contributions — phantom, validated simulator, RL.",
    full_bullets=[
        (0, "Adaptive MRI: choose the next acquisition from the current estimate (≈2.5× faster, Beracha)", False),
        (0, "The tension:", True),
        (1, "adaptive qMRI exists — but only as a hand-derived BAYESIAN model", False),
        (1, "RL has controlled MRI scanners — but for NON-quantitative goals (k-space, shape)", False),
        (1, "nobody has put a LEARNED policy on the QUANTITATIVE objective", True),
        (0, "Thesis: the first RL agent for adaptive quantitative MRI", True),
        (1, "conditions each acquisition on the current fitted T1; trained & scored on T1 error", False),
        (0, "Roadmap: trustworthy phantom (A1) → validated simulator (A2) → the RL (A3)", True)])

# SLIDE 4 — A1 twin
content_slide(
    "A1 — A known phantom to learn from",
    "Execution (50%) — the digital twin + reproducibility signal",
    [(0, "Need a KNOWN object: ground-truth T1/T2", True),
     (0, "QalibreMD Model 130 — 14 spheres, ~24 ms–1.9 s", False),
     (0, "MRISystemPhantom.jl → KomaMRI objects", True),
     (0, "Water (~80% spins) on its own coarsenable grid", True),
     (0, "Open-source, ~373 tests, config + seed reproducible", False)],
    ["potential_calibremd_phantom_replacement.png"],
    "phantom figure",
    "A1. The agent learns entirely in simulation so it needs a known object. QalibreMD "
    "Model 130, the standard NIST/ISMRM phantom, 14 spheres ~24ms-1.9s. I built "
    "MRISystemPhantom.jl: configurable open-source twin returning KomaMRI objects, plus "
    "ground truth, fitter, randomisation. Design decision: water (80% of spins) on its own "
    "coarsenable grid — the lever for cheap RL later. Reproducible: tested package, "
    "serialisable PhantomConfig+seed, any result regenerates from a config.",
    full_bullets=[
        (0, "Agent learns in simulation → needs a KNOWN object with ground-truth T1/T2", True),
        (0, "QalibreMD Model 130: standard NIST/ISMRM calibration phantom", False),
        (1, "14 spheres spanning the clinical T1 range (~24 ms to 1.9 s)", False),
        (0, "Built MRISystemPhantom.jl — configurable twin returning KomaMRI objects", True),
        (1, "ground-truth labels, the T1 fitter, per-episode randomisation", False),
        (0, "Key decision: background water (~80% of spins) on its OWN coarsenable grid", True),
        (1, "the lever that later makes RL training affordable", False),
        (0, "Reproducible by construction: open-source pkg, ~373 tests, serialisable config + seed", True)])

# SLIDE 5 — Live demo
content_slide(
    "LIVE DEMO — Building and inspecting the digital phantom",
    "Communication / Execution — show the twin is real, runnable engineering",
    [(0, "Build phantom live from a PhantomConfig", True),
     (0, "3-D T1 plate → Bloch sim → one IR signal", False),
     (0, "A real KomaMRI.Phantom — no special-casing", False),
     (0, "Fallback stills · pre-warm JIT · hard stop 5:30", True)],
    ["phantomdocs.png"],
    "demo fallback still",
    "DEMO. Run examples/plot_phantom_3d.jl / t1_mapping.jl live: build from a config, show "
    "the 3-D voxelised plate coloured by T1, simulate one IR signal. Narrate that it's a "
    "real KomaMRI.Phantom. Fallback screenshots on the slide. Pre-warm Julia JIT. Hard stop "
    "at 5:30.",
    full_bullets=[
        (0, "Build a phantom live from a PhantomConfig", True),
        (0, "Show the 3-D voxelised T1 plate, coloured by T1 value", False),
        (0, "Hand it straight to the Bloch simulator → one IR signal / reconstructed image", False),
        (0, "“This is a real KomaMRI.Phantom — no special casing.”", False),
        (1, "FALLBACK: pre-captured screenshots on this slide if the build stalls", True),
        (1, "pre-warm the Julia session before the talk (avoid first-call JIT)", False),
        (1, "HARD STOP at 5:30 — do not overrun", True)])

# SLIDE 6 — A2 plausible != correct
content_slide(
    "A2 — A plausible image is not a correct measurement",
    "Execution (50%) — validation rigour: parameter recovery, not eyeballing",
    [(0, "Simulator must be QUANTITATIVELY right", True),
     (0, "KomaMRI validated only on SHORT sequences", False),
     (0, "Recovery test, no noise: 39.4% error (some ~100%)", True),
     (0, "Spoilers worse, long TR worse → elapsed TIME, not physics", True)],
    ["komaMRI/buggy_t1_fit_vs_true.png"],
    "buggy T1 fit",
    "A2. A simulator the agent learns from must be quantitatively correct, not just "
    "plausible. KomaMRI benchmarked only on short sequences. Validate by parameter "
    "recovery, no noise: should be exact, got 39.4% error, images looked fine. Hard part "
    "was diagnostic: spoilers made it worse (ruling out the obvious physics explanation), "
    "longer TR worse too — pointing at elapsed time, not physics.",
    full_bullets=[
        (0, "If the agent learns from a simulator, it must be QUANTITATIVELY right", True),
        (1, "KomaMRI is peer-reviewed — but benchmarked only on SHORT sequences", False),
        (0, "Validate as qMRI demands: known phantom → fit T1 → compare, no noise", True),
        (0, "It should be near-exact. It wasn't: 39.4% mean error, some near 100%", True),
        (1, "and the images still looked fine", False),
        (0, "Ruling out physics (the hard part):", True),
        (1, "adding spoilers made it WORSE (39.4 → 44.0%)", False),
        (1, "longer TR made it worse → driven by total elapsed TIME, not physics", False)])

# SLIDE 7 — A2 the bug
content_slide(
    "A2 — Two floating-point bugs in KomaMRI, found & fixed",
    "Execution (50%) — minimal reproducers, upstream contribution, the 'aha'",
    [(0, "k-space collapses after a fixed cumulative TIME", True),
     (0, "RF edge marked with a fixed ε = 1e-14 s", True),
     (0, "Float64 is relative: eps(t) ≈ t·2⁻⁵² → t+ε == t past ~128 s", False),
     (0, "Fixed upstream (PR #780, #789): 39.4% → 0.48%", True)],
    ["fig3_float_collapse.png", "komaMRI/buggy_pixel_grid_overlay.png"],
    "FIG-3 float-collapse graph\n+ buggy pixel-grid overlay",
    "A2 bug. Diagnostic: look at k-space directly — collapses after fixed cumulative time, "
    "regardless of shot count. Mechanism: fixed absolute ε=1e-14 marks the RF edge, but "
    "Float64 precision is relative (eps(t)≈t·2⁻⁵²); past ~128s t+ε==t, edge markers "
    "collapse, wrong RF for one sample = 5% of a 1ms pulse → 20-25% jump → biased T1. "
    "Minimal reproducers, fixed upstream, PR #780 merged, #789 open. 39.4→0.48%. This is "
    "what makes Part 3 trustworthy.",
    full_bullets=[
        (0, "Decisive diagnostic: k-space collapses after a fixed cumulative TIME", True),
        (0, "Mechanism: RF edge marked with a FIXED ε = 1e-14 s…", True),
        (1, "…but Float64 precision is RELATIVE: gap grows as eps(t) ≈ t·2⁻⁵²", False),
        (1, "past ~128 s, t+ε == t bit-for-bit → edge markers collapse", False),
        (1, "wrong RF for one sample = 5% of a 1 ms pulse → 20–25% signal jump", False),
        (0, "Reduced to minimal reproducers, fixed, contributed upstream: PR #780 (merged), #789", True),
        (0, "Result: 39.4% → 0.48% mean error (max 1.2%)", True),
        (1, "novel failure mode for an MRI simulator; bites exactly in the long-sequence RL regime", False)])

# SLIDE 8 — RL formulation
content_slide(
    "A3 — The RL formulation",
    "Execution (50%) — the MDP and the engineering of the loop",
    [(0, "Episode: 2-D IR scan, 240 s budget", True),
     (0, "State: per-sphere T1 est + budget", True),
     (0, "Action: next TI (and TR)", True),
     (0, "Reward: dense Δlog MAPE", True),
     (0, "PPO (SB3), Julia in-process", False),
     (0, "σ=50 → SNR 17–28", False)],
    ["fig_envloop.png"],
    "env-loop diagram\n(export TikZ fig\nfrom the report)",
    "A3 formulation. Episode = 2-D IR scan through the T1 plate under a 240s budget. State: "
    "running per-sphere T1 estimates + budget; agent sees estimates not pixels. Action: next "
    "TI (and TR). Reward dense: Δlog MAPE — because an early terminal-bonus version "
    "collapsed to a fixed policy. PPO via SB3, Julia in-process. Noise σ=50 → SNR 17-28, "
    "clinical.",
    text_w=3.7,
    full_bullets=[
        (0, "Episode = a 2-D inversion-recovery scan through the T1 plate, fixed time budget (240 s)", True),
        (0, "State: running per-sphere T1 estimates + budget used", True),
        (1, "agent never sees pixels — only the fitter's estimates (real qMRI interface)", False),
        (0, "Action: the next block's timings — inversion time TI (and TR)", True),
        (0, "Reward: DENSE — per-step improvement in log-error, Δlog MAPE", True),
        (1, "why dense: an early version with a big terminal bonus collapsed to a fixed policy", False),
        (0, "Stack: Gymnasium + Stable-Baselines3 PPO, Python↔Julia, simulator in-process", False),
        (0, "Noise calibrated: σ=50 → SNR ≈ 17–28, clinical IR range", False)])

# SLIDE 9 — multi-fidelity
content_slide(
    "A3 / C2 — The cost wall and the multi-fidelity ladder",
    "Execution (50%) — adapting a technique with INSIGHT (cached-water linearity)",
    [(0, "Full-Bloch step ≈ 4 s → 200k steps ≈ 9 days", True),
     (0, "80% of spins are water the agent never measures", True),
     (0, "Lever: coarsen water · cache water", True),
     (0, "~30× cost ladder, one policy warm-started up", True)],
    ["fig4_ladder.png"],
    "FIG-4 ladder\nstaircase: analytic→cached3→\ncached→full3→full, cost up,\nwater-spins bar shrinking",
    "C2 cost wall. Full-Bloch step ~4s CPU → 200k steps ~9 days. 80% of spins are water the "
    "agent doesn't measure. Two levers: coarsen water (rescale PD to conserve "
    "magnetisation), or cache it exploiting linearity S_full=S_spheres+S_water — homogeneous "
    "water factorises, new timing just rescales a template computed once, ~8x saving, matches "
    "full-Bloch to 0.12%. Net ~30x ladder, warm-start one policy up the rungs. Next slide: the "
    "linearity that makes caching exact.",
    full_bullets=[
        (0, "The wall: a full-Bloch step ≈ 4 s on CPU → a 200k-step run is ~9 DAYS", True),
        (1, "RL normally trains on environments millions of times cheaper", False),
        (0, "Key observation: 80% of spins are background water the agent never measures", True),
        (0, "Two levers, both on the water:", True),
        (1, "coarsen it (own grid, proton density rescaled to conserve magnetisation)", False),
        (1, "cache it — exploit simulator LINEARITY: S_full = S_spheres + S_water", False),
        (2, "homogeneous water factorises → new timing just rescales a cached template", False),
        (1, "~8× per-step saving, matching full-Bloch T1 fits to 0.12%", False),
        (0, "Net: a ~30× cost-range ladder, warm-starting one policy up the rungs", True)])

# SLIDE 9b — cached-water linearity (equations only)
content_slide(
    "A3 / C2 — Cached water: exploiting simulator linearity",
    "Execution (50%) — the insight that makes the cheap rungs trustworthy",
    [(0, "Bloch is LINEAR in spins → k-space separates", True),
     (0, "Homogeneous water FACTORISES per shot", True),
     (0, "Build template once, rescale per timing", False),
     (0, "~8× cheaper · 0.12% T1 match", True)],
    ["fig4b_linearity.png"],
    "FIG-4b equations\nS_full = S_spheres + S_water\nS_water[k] = Mz·sinα·W_α[k]",
    "Cached-water linearity. Bloch is linear in spins so k-space separates exactly: "
    "S_full = S_spheres + S_water. Water is one homogeneous material so its contribution "
    "factorises per shot — only the scalar Mz^(k) depends on the timing, the geometric "
    "template W_α is built once and reused. A new timing just rescales the template; only the "
    "sphere spins are re-simulated. ~8× cheaper, matches full-Bloch to 0.12%. This is what "
    "makes the cheap rungs trustworthy, not just fast.",
    text_w=4.5,
    full_bullets=[
        (0, "Bloch simulation is LINEAR in spins → the k-space separates exactly", True),
        (0, "Water is a single homogeneous material → its term FACTORISES per shot", True),
        (1, "only the scalar Mz^(k) depends on the timing (TI, TR, α)", False),
        (1, "the geometric template W_α is built ONCE, then reused", False),
        (0, "So a new timing just rescales a cached template — only spheres re-simulated", True),
        (1, "~8× cheaper per step; matches full-Bloch T1 fits to 0.12%", False)])

# SLIDE 10 — promotion
content_slide(
    "A3 — The trust problem: bias-aware promotion",
    "Execution (50%) / Evaluation (20%) — methodological care",
    [(0, "Cheap biased rung games its own score", True),
     (0, "Score on a held-out FULL-sim probe", True),
     (0, "Promote on plateau / ranking breakdown", False),
     (0, "Keep a GLOBAL-BEST checkpoint", True)],
    ["fig5_controller.png"],
    "controller diagram",
    "Promotion. A cheap biased rung games its own score, so promotion and checkpointing are "
    "scored on a held-out full-sim probe, never the cheap number. Promote on plateau or "
    "ranking breakdown. Keep a global-best checkpoint on target fidelity — the last policy "
    "is often not the best. Honest: thresholds hand-tuned, future work.",
    text_w=4.5,
    full_bullets=[
        (0, "A cheap, biased rung can keep RAISING ITS OWN SCORE by exploiting its bias", True),
        (0, "So promotion + checkpoint selection are scored on a held-out FULL-simulator probe", True),
        (1, "never the cheap rung's own number", False),
        (0, "Promote on: target plateau, or ranking breakdown (Spearman vs full sim drops)", False),
        (0, "Keep a GLOBAL-BEST checkpoint on the target fidelity", True),
        (1, "the last policy in a curriculum is often NOT the best (we'll see this in Run A)", False),
        (0, "Honest note (Q&A): thresholds are hand-tuned, not derived — flagged as future work", False)])

# SLIDE 11 — Run A limitation
content_slide(
    "Limitation — where adaptivity does not help (Run A)",
    "Evaluation (20%) — the honest negative result, dedicated framing",
    [(0, "Run A: full 14-sphere plate, whole T1 range", True),
     (0, "Fixed Cramér–Rao WINS: 4.70% vs 9.4–11.7%", True),
     (0, "Agent is adaptive (TI tracks estimate), not collapsed", False),
     (0, "Longest sphere beyond TI reach — coverage ceiling", False),
     (0, "Lessons → global-best; test where there's variation", True)],
    ["e2_rl/runA_moneyplot.png", "e2_rl/runA_ti_per_episode.png"],
    "Run A behaviour",
    "Limitation slide. Run A = full 14-sphere plate spanning the whole T1 range — hardest "
    "case for adaptivity, one fixed schedule already near-optimal. Fixed CR wins 4.70% vs "
    "9.4-11.7%. But the agent IS adaptive (TI varies, tracks estimate); longest sphere "
    "beyond TI reach. Lessons: global-best necessary; need a task with variation. Motivates "
    "Run B.",
    full_bullets=[
        (0, "Run A = full 14-sphere plate, spanning the ENTIRE clinical T1 range", True),
        (1, "the hardest case for adaptivity — one fixed schedule already serves a near-fixed fleet", False),
        (0, "Result: the fixed Cramér–Rao schedule WINS (4.70% vs the agent's 9.4–11.7%)", True),
        (1, "an honest negative", False),
        (0, "But the agent is genuinely adaptive, not collapsed:", True),
        (1, "TI varies within/across episodes, tracks the running T1 estimate", False),
        (1, "longest sphere is beyond its TI reach — an action-coverage ceiling", False),
        (0, "Lessons → Run B: global-best is necessary; test adaptivity where there's something to adapt to", True)])

# SLIDE 12 — Run B
content_slide(
    "Run B — adaptivity helps when the task varies",
    "Execution + Evaluation — fair comparison, alternative explanation ruled out",
    [(0, "Fair: held-out seeds, 24 eps, CIs, matched budget + fitter", False),
     (0, "5 spheres, T1 resampled every episode", True),
     (0, "RL 4.62% vs fixed 6.86% (CIs disjoint); CR 15.96%", True),
     (0, "Not just more blocks: regrid at policy TR → worse", True),
     (0, "Adaptive PLACEMENT; holds at 560 s (4.16 vs 6.04)", False)],
    ["e2_rl/runB_mape_comparison.png"],
    "Run B comparison",
    "Run B. Say the fair-comparison line once: held-out seeds, 24 episodes, CIs, matched "
    "budget + same fitter. 5 spheres, T1 resampled each episode. RL 4.62% vs fixed 6.86%, "
    "CIs don't overlap; CR collapses (15.96%). Is it just more measurements? No — re-solve "
    "grid at policy's own TR, it does worse → gain is adaptive placement. Probe-then-refine; "
    "holds at 560s.",
    full_bullets=[
        (0, "Fair comparison (covers every result): strict held-out seeds, 24 episodes, 95% CIs,", True),
        (1, "same scan budget + same fitter for RL and every fixed baseline", False),
        (0, "Run B: 5 spheres, T1 resampled continuously every episode (variation a grid can't anticipate)", True),
        (0, "RL 4.62% vs best fixed log-grid 6.86% at 240 s — non-overlapping CIs", True),
        (1, "Cramér–Rao schedule collapses (15.96%): solved for a nominal fleet", False),
        (0, "Is the gain just MORE measurements? No.", True),
        (1, "re-solve the fixed grid at the policy's own shorter TR → it does WORSE", False),
        (1, "→ the gain is adaptive PLACEMENT, not block count", True),
        (0, "Behaviour: probe-then-refine. Advantage holds at 560 s (4.16% vs 6.04%)", False)])

# SLIDE 13 — memory ablation (climax)
content_slide(
    "Key finding — what the policy should remember",
    "Execution + Evaluation — a genuine research insight, controlled ablation",
    [(0, "History = which TIs you've bought (partially observed)", True),
     (0, "Carry it: σ · LSTM · TI-coverage histogram", False),
     (0, "histogram ≫ none > σ > LSTM", True),
     (0, "Best: 2.93% MAPE, 95.8% < 5% (vs 6.04%)", True),
     (0, "LSTM tracked estimate (r=+0.78) & lost; histogram coverage (r=−0.12) & won", False)],
    ["e2_rl/runB_memory_mape.png"],
    "memory ablation (money result)",
    "Climax. Partially-observed problem: estimate is a lossy summary of which TIs bought. "
    "Three memories: σ, LSTM, TI-coverage histogram. histogram >> no-mem > σ > LSTM. "
    "Histogram 2.93%, 95.8% under 5%, vs 6.04% fixed, CI separate from all. Twist: LSTM "
    "learned the intuitive estimate-tracking rule (r=+0.78) and lost; histogram is "
    "coverage-driven (r=-0.12) and won. What you remember matters as much as whether you "
    "adapt. LSTM caveat: 5x slower, undertrained.",
    full_bullets=[
        (0, "Partially-observed: the estimate is a lossy summary of WHICH TIs you've bought", True),
        (0, "Three ways to carry that history: fitter uncertainty (σ), an LSTM, a TI-coverage histogram", False),
        (0, "Result: histogram ≫ no-memory > σ-channel > LSTM", True),
        (1, "histogram = best in the project: 2.93% MAPE, 95.8% of episodes < 5% (vs 6.04% fixed)", True),
        (0, "The twist:", True),
        (1, "LSTM learned the INTUITIVE rule (longer T1 → longer TI, r=+0.78) — and LOST", False),
        (1, "histogram conditions on which TI bins are unsampled (r=−0.12) — and WON", False),
        (0, "Takeaway: what you remember matters as much as whether you adapt", True),
        (1, "LSTM caveat (Q&A): ~5× slower/step → undertrained at equal wall-clock", False)])

# SLIDE 14 — positioning
content_slide(
    "Quantitative positioning against published work",
    "Evaluation (20%) — contribution situated vs baselines, honest caveats",
    [(0, "Not like-for-like — tasks differ", True),
     (0, "Prior work: stronger hardware validation", False),
     (0, "Ours: closed-loop RL on fitted-T1, validated pipeline", True),
     (0, "Beracha heuristic: 6.04→2.93% ≈ ~4× acceleration", True)],
    ["FIG9_anchors.png"],
    "FIG-9 (make)\ndistil quant-anchor table\nto 3–4 comparable bars",
    "Positioning. Not like-for-like; prior work has stronger hardware validation. We "
    "contribute a different object: closed-loop RL scored on fitted-T1 error in a validated "
    "pipeline. By Beracha's own precision->time heuristic, 6.04->2.93% (~2.06x precision) ~ "
    "4x acceleration, top of their range, harder task.",
    full_bullets=[
        (0, "Not a like-for-like ranking — tasks differ", True),
        (0, "Prior work has stronger hardware validation (Beracha, MRzero, Walker-Samuel)", False),
        (0, "This work contributes a different OBJECT:", True),
        (1, "a closed-loop RL policy scored on fitted-T1 error in a VALIDATED 2-D pipeline", False),
        (0, "Using Beracha's own precision→time heuristic:", True),
        (1, "6.04% → 2.93% (~2.06× precision) ≈ ~4× acceleration", False),
        (1, "top of their reported range — on a harder fleet task, not one voxel", True)])

# SLIDE 15 — close
content_slide(
    "Contributions & future work",
    "Evaluation (20%) — broader impact + what does NOT generalise",
    [(0, "1. Open-source phantom twin", True),
     (0, "2. Validation → fixed two KomaMRI bugs (39.4→0.48%)", True),
     (0, "3. First RL agent for adaptive qMRI (2.93 vs 6.04)", True),
     (0, "Impact: ~4× acceleration, MR-Linac bottleneck", True),
     (0, "Not claiming: sim-only, one seed, fitter confound", False),
     (0, "Future: wider actions; T2/PD; Bayesian baseline; LSTM", False)],
    [], "", has_img=False,
    notes=(
        "Close. Three contributions: twin; validation fixing two upstream bugs; first RL agent "
        "for adaptive qMRI (2.93 vs 6.04, plus coverage>estimate-tracking finding). Why it "
        "matters: scan-time bottleneck, MR-Linac, ~4x accel. What I'm NOT claiming: sim-only, "
        "one seed, single-fitter confound; next step physical phantom; 14-sphere favours fixed. "
        "Future work. Thanks — questions."),
    full_bullets=[
        (0, "1. Open-source executable twin of a calibrated phantom", True),
        (0, "2. Validation-by-recovery → fixed two upstream KomaMRI bugs (39.4 → 0.48%)", True),
        (0, "3. First RL agent for adaptive qMRI — best 2.93% vs 6.04%; coverage-memory finding", True),
        (0, "Why it matters: less scan time at equal accuracy = the MR-Linac bottleneck (~4× accel)", True),
        (0, "What I'm NOT claiming (reflection): simulator-only, one seed/arm, single-fitter confound", True),
        (1, "honest next step is a physical-phantom benchmark; full 14-sphere still favours fixed", False),
        (0, "Future: wider action space; T2/PD plates; greedy-CRLB Bayesian baseline; step-matched LSTM", False)])

# ---------------------------------------------------------------------------
# BACKUP DIVIDER
# ---------------------------------------------------------------------------
s = add_slide()
band = s.shapes.add_shape(1, Inches(0), Inches(3.0), SW, Inches(1.5))
band.fill.solid(); band.fill.fore_color.rgb = NAVY
band.line.fill.background(); band.shadow.inherit = False
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.25), Inches(11.7), Inches(1.0))
p = tb.text_frame.paragraphs[0]
p.text = "Backup slides  —  anticipating Q&A"
p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)
set_notes(s, "Backup deck = anticipating objections (an Execution-band signal). "
             "Likely challenges: fitter confound, single seed, simulator-only, LSTM compute caveat.")


def backup_slide(title, bullets, images, img_label, notes):
    return content_slide(title, "Backup — Q&A", bullets, images, img_label, notes)


backup_slide(
    "Backup — Reward screening (why Δlog MAPE)",
    [(0, "Per-step −MAPE penalty collapsed to a ZERO-measurement policy (summed penalties punish measuring)", False),
     (0, "Terminal-only was too sparse; progress rewards best", False),
     (0, "Δlog MAPE won the 30k PPO screen (0.830) — chosen before committing Bloch compute", False),
     (0, "Ties back to the degenerate single-voxel collapse (terminal-bonus lesson)", False)],
    ["FIG_reward.png"], "reward-screen table (tab:e2-reward-screen)",
    "Reward screening detail. neg-MAPE collapses to no blocks; terminal too sparse; "
    "Δlog MAPE best (0.830). Chosen on the cheap analytic env before Bloch compute.")

backup_slide(
    "Backup — Cached-water validation",
    [(0, "Cached model matches full-Bloch T1 fits to 0.12% (worst 1.15%)", False),
     (0, "Modelling water analytically instead is too crude: up to ~40% per-sphere deviation", False),
     (0, "α-matched template bank: exact at matched flip angle, blended off-grid", False)],
    ["e2_rl/water_t1_fit_4variants.png", "e2_rl/water_cache_relerr_vs_TI.png"],
    "cached-water validation",
    "Cached water matches full-Bloch to 0.12%; analytic water ~40% off. Template bank keyed "
    "to flip angle.")

backup_slide(
    "Backup — Noise calibration (σ = 50)",
    [(0, "σ=50 → NEMA SNR ≈ 17 (plate avg) to 28 (bright spheres) — clinical IR range", False),
     (0, "25-seed sweep: T1 fits degrade gracefully to σ≈85, collapse beyond", False),
     (0, "So the task is noise-LIMITED but not estimation-broken", False)],
    ["e2_rl/snr_mape_vs_sigma.png", "e2_rl/snr_sweep_images.png"],
    "SNR sweep",
    "Noise σ=50 → SNR 17-28, clinical. Graceful to σ≈85. Noise-limited, not broken.")

backup_slide(
    "Backup — The second KomaMRI bug + fitter floor",
    [(0, "Bug 2: closing-knot collapse AFTER time-rebasing (T0 + t) — same ε mechanism, later in pipeline", False),
     (1, "reproducer: +52% jump at shot 18 (270 s), then a wrong plateau", False),
     (0, "Fitter grid floor: 300 log-spaced T1 candidates → ~1% quantisation floor by construction", False),
     (1, "so reported MAPE can't fall below ~1%; 500-point grid reaches ~0.48%", False)],
    ["komaMRI/fixed_recovery_curves_koma.png"], "fixed recovery curves",
    "Second bug: closing-knot collapse after rebasing, 52% jump at shot 18. Fitter floor: "
    "300-point log grid → ~1% quantisation; 500 points ~0.48%.")

backup_slide(
    "Backup — Run B policy behaviour",
    [(0, "240 s policy: 5.92 blocks avg, ZERO timing repairs", False),
     (0, "Opens with two high-TI probes (~0.78–0.83 s), then shorter adaptive refinements", False),
     (0, "TI vs current estimate r=−0.72: a two-phase probe-then-refine, not a monotonic rule", False)],
    ["e2_rl/runB_240_ti_per_episode.png", "e2_rl/runB_curriculum_probes.png"],
    "Run B policy",
    "Run B behaviour: ~6 blocks, zero repairs, two high-TI probes then refine, r=-0.72 "
    "two-phase strategy.")

backup_slide(
    "Backup — Fixed comparators (Cramér–Rao baselines)",
    [(0, "Fixed schedules aren't hand-chosen — they're Cramér–Rao-optimal", False),
     (0, "CRLB = variance floor for any unbiased estimator, set by design + noise", False),
     (0, "static-OPTIMAL (solved for nominal fleet) vs static-AVERAGE (robust log grid)", False),
     (1, "CR wins when the fleet matches (Run A); collapses under per-episode variation (Run B)", False)],
    ["FIG_crlb.png"], "(no figure — speak it)",
    "Fixed comparators are CR-optimal, not hand-picked. CRLB = variance floor. "
    "static-optimal vs static-average; CR wins Run A, collapses Run B.")

prs.save(OUT)
print("Wrote", OUT, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
